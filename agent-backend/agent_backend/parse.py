"""File-to-text extraction for source materials.

Each parser takes raw bytes and returns plain text. The dispatcher picks the
parser by file extension. Errors surface as ParseError with a friendly message.
"""
from __future__ import annotations

import io
import logging
import os
from typing import Callable

log = logging.getLogger(__name__)


class ParseError(Exception):
    pass


def _parse_txt(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("utf-8", errors="replace")


def _parse_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ParseError("pypdf not installed") from exc
    try:
        reader = PdfReader(io.BytesIO(data))
        pages = [p.extract_text() or "" for p in reader.pages]
        return "\n\n".join(pages).strip()
    except Exception as exc:
        raise ParseError(f"could not read PDF: {exc}") from exc


def _parse_pptx(data: bytes) -> str:
    """Backwards-compatible flat-text PPTX parse.

    Kept for callers that just want the joined text (e.g. simple
    grounding into a generic Anthropic prompt). For richer structured
    parses, use _parse_pptx_structured() below — it returns slide-level
    metadata so the agent can map module / lesson boundaries to slide
    ranges and cite slides in source attribution.
    """
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ParseError("python-pptx not installed") from exc
    try:
        prs = Presentation(io.BytesIO(data))
        slides_out: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            chunks = [f"Slide {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            chunks.append(text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    chunks.append(f"(Notes) {notes}")
            slides_out.append("\n".join(chunks))
        return "\n\n".join(slides_out).strip()
    except Exception as exc:
        raise ParseError(f"could not read PPTX: {exc}") from exc


# ─── Track-SD (Source-Deck deepen): structured PPTX parse ─────────────────
#
# The flat-text parse above gives the agent one big concatenated string —
# fine for short decks, but for 30+ slide decks the agent loses track of:
#   1. WHICH slide a piece of content came from (no citation possible)
#   2. WHERE natural module breaks fall (section dividers get blended in)
#   3. WHICH slides became which lesson (no slide-range attribution)
#
# Structured parse returns:
#   {
#     slides: [{ n, title, body, notes, isSection }, ...],
#     totalSlides: int,
#     sectionCount: int,
#   }
#
# The agent's MODE 1 (Course Architect) receives this alongside the flat
# text and uses it to detect natural module boundaries (section slides),
# cite slide ranges per lesson, and propose a clean slide → module map
# the LD can verify before generation.

def _parse_pptx_structured(data: bytes) -> dict:
    """Return structured slide-level metadata for a PPTX.

    Each slide entry has:
      n         1-indexed slide number
      title     first non-empty text frame (heuristic: top of slide,
                largest font usually). Empty if no text frames.
      body      remaining text frames concatenated, newline-separated
      notes     speaker notes text (empty if none)
      isSection True if this slide looks like a section divider —
                heuristic: short title (<= 8 words), no body, not
                slide 1 (which is the deck title not a section).
    """
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ParseError("python-pptx not installed") from exc
    try:
        prs = Presentation(io.BytesIO(data))
        slides: list[dict] = []
        for i, slide in enumerate(prs.slides, start=1):
            text_frames: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # Concat all paragraphs in this text frame, preserving
                # newlines so bullets stay readable.
                lines: list[str] = []
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)
                if lines:
                    text_frames.append("\n".join(lines))

            title = text_frames[0] if text_frames else ""
            body = "\n\n".join(text_frames[1:]) if len(text_frames) > 1 else ""

            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            # Section heuristic: short title + no body + not slide 1.
            # Title-only slides with brief headers are almost always
            # section dividers in BCG-style decks.
            title_word_count = len(title.split())
            is_section = (
                i > 1
                and title != ""
                and not body
                and title_word_count > 0
                and title_word_count <= 8
            )

            slides.append({
                "n": i,
                "title": title,
                "body": body,
                "notes": notes,
                "isSection": is_section,
            })

        section_count = sum(1 for s in slides if s["isSection"])
        return {
            "slides": slides,
            "totalSlides": len(slides),
            "sectionCount": section_count,
        }
    except Exception as exc:
        raise ParseError(f"could not parse PPTX structure: {exc}") from exc


def _parse_docx(data: bytes) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise ParseError("python-docx not installed") from exc
    try:
        doc = Document(io.BytesIO(data))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip(" |"):
                    parts.append(row_text)
        return "\n".join(parts).strip()
    except Exception as exc:
        raise ParseError(f"could not read DOCX: {exc}") from exc


_PARSERS: dict[str, Callable[[bytes], str]] = {
    ".txt": _parse_txt,
    ".md": _parse_txt,
    ".markdown": _parse_txt,
    ".pdf": _parse_pdf,
    ".pptx": _parse_pptx,
    ".docx": _parse_docx,
}

SUPPORTED_EXTENSIONS = sorted(_PARSERS.keys())


def parse_file(filename: str, data: bytes) -> str:
    """Extract text from a file by its extension. Raises ParseError on failure."""
    ext = os.path.splitext(filename.lower())[1]
    parser = _PARSERS.get(ext)
    if parser is None:
        raise ParseError(
            f"unsupported file type '{ext}' — supported: {', '.join(SUPPORTED_EXTENSIONS)}"
        )
    return parser(data)
